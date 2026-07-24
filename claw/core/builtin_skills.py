"""Built-in ("pre-built") skills shipped with the platform.

Unlike user skills (rows in the ``skills`` table), these are defined in code so
they're always available, version-controlled, and can't be edited or deleted.
They're merged into the agent's skills list — the system-prompt summary and the
``read_skill`` tool — and surfaced read-only in Settings → Skills.

``skill-creator`` teaches the agent how to author new user skills correctly via
the ``manage_skill`` tool (writing workspace files does not create a skill,
which is the trap that made agent-"created" skills never show up). The
``pptx``/``xlsx``/``pdf``/``docx`` document skills are recipes for the exec
tool: they teach techniques for the document/data stack already baked into the
sandbox image (docker/sandbox.Dockerfile) — python-pptx, openpyxl, PyMuPDF,
python-docx, reportlab, pandas, etc. — rather than introducing new in-process
tools, since document creation/editing already happens by shelling out to the
sandbox.
"""

from dataclasses import dataclass, field
from datetime import datetime, timezone

# Stable timestamp so the API payload is deterministic across restarts.
_BUILTIN_TS = datetime(2026, 1, 1, tzinfo=timezone.utc)


@dataclass(frozen=True)
class BuiltinSkill:
    """A code-defined skill; duck-types the ORM ``Skill`` for the summary/API."""

    name: str
    description: str
    content: str
    enabled: bool = True
    # Optional "capabilities covered" — shown as a 2x2 card grid in the
    # Settings → Skills detail view. Empty for skills with no such UI (e.g.
    # skill-creator), which fall back to the plain raw-content view instead.
    capabilities: tuple[tuple[str, str], ...] = field(default_factory=tuple)
    summary: str = ""

    @property
    def id(self) -> str:
        return f"builtin:{self.name}"

    @property
    def updated_at(self) -> datetime:
        return _BUILTIN_TS


_SKILL_CREATOR_CONTENT = """\
A skill is a reusable procedure stored in the system and offered to you in every
future chat. A skill is NOT a workspace file: writing a `.py`/`.md` file does NOT
create a skill and it will never appear in Settings → Skills. The ONLY way to
create or update a skill is the `manage_skill` tool.

## Create a skill (manage_skill, action="save")
1. name — short, kebab-case, unique per user (e.g. `stock-analysis`, `weekly-report`).
2. description — ONE line, shown to you in every chat's system prompt. This is how
   future-you decides to use the skill, so make it specific and trigger-oriented:
   what it does AND when to use it. e.g. "Analyse a Thai/US stock over N months
   (price trend, dividend yield, PE) — use when the user asks to analyse a stock."
3. content — the full instructions, loaded on demand when you `read_skill`. Keep it
   procedural and focused:
   - The exact steps, in order.
   - Which tools to use at each step (web_search, mcp_* connectors, exec for
     computation/charts, write_file for outputs).
   - The output the user expects (table, chart PNG, PDF) and its format.
   - Sensible defaults and edge cases (default period, currency, data source).
4. Call `manage_skill` with action="save" and those fields.

## What makes a good skill
- One skill = one clear procedure. Don't bundle unrelated tasks.
- Capture the reusable *method*, not one-off data from this chat.
- Name the tools you'll call so future-you needs no guesswork.
- Keep content concise — you pay tokens to read it; link steps, don't pad.
- After saving, tell the user it's saved and now appears in Settings → Skills
  (enabled by default).

## Manage existing skills
- action="list" — list your saved skills.
- action="save" — create, or overwrite one with the same name.
- action="delete" — remove a skill by name.

Never try to "install" a skill by writing files or editing the database — always
go through `manage_skill`.
"""

_PPTX_CONTENT = """\
Build and edit PowerPoint decks with `python-pptx` (already installed in the sandbox)
via the `exec` tool. Write the .pptx to /workspace so it's downloadable when done.

## Presentation Creation
For a new deck, don't hand-place shapes by guessing coordinates — design the slide as
an HTML/CSS layout first (absolute-positioned divs at the deck's exact pixel size,
e.g. 1280x720 for 16:9), render it to a PNG with the tools already in the sandbox
(`weasyprint` HTML→PDF, then PyMuPDF `fitz` to rasterize page 0 to PNG), and look at
the image before committing to it. This catches overlap/overflow far faster than
reasoning about EMU coordinates blind. Once the layout is right, translate each div's
`(left, top, width, height)` in pixels to inches (`px / 96`) and build the real shapes:

```python
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 16:9
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout — full control

PALETTE = {"primary": RGBColor(0x1F, 0x4E, 0x79), "accent": RGBColor(0xED, 0x7D, 0x31)}

box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
p = box.text_frame.paragraphs[0]
p.text = "Q3 Supplier Review"
p.font.size, p.font.bold, p.font.color.rgb = Pt(32), True, PALETTE["primary"]

table = slide.shapes.add_table(rows=4, cols=3, left=Inches(0.5), top=Inches(1.6),
                                width=Inches(9), height=Inches(2.5)).table

chart_data = CategoryChartData()
chart_data.categories = ["Q1", "Q2", "Q3"]
chart_data.add_series("Revenue", (19.2, 21.4, 25.0))
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(4.3),
                        Inches(9), Inches(2.8), chart_data)
prs.save("/workspace/review.pptx")
```

## Template Reuse
To reuse an existing deck as a template (matching layouts, duplicating a slide, then
replacing its content) rather than building from scratch, `python-pptx` has no native
"duplicate slide" API — build it via `copy.deepcopy` on the slide's XML:

```python
import copy
from pptx.oxml.ns import qn

def duplicate_slide(prs, index):
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    # A layout-based add_slide() pre-populates inherited placeholders — strip them
    # before copying the source's own shapes over, or they collide/duplicate.
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in source.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shape._element))
    return new_slide

new_slide = duplicate_slide(prs, 0)
for shape in new_slide.shapes:
    if shape.has_text_frame and "Q3" in shape.text_frame.text:
        shape.text_frame.paragraphs[0].runs[0].text = "Q4 Supplier Review"
```
Match layouts by reusing `source.slide_layout` (not a fresh blank layout) so fonts,
placeholder positions, and theme colors carry over automatically.

## Slide Editing
For edits `python-pptx`'s object model doesn't expose directly (custom XML attributes,
some line/fill effects), drop to its OOXML tree via `shape._element` /
`slide.shapes._spTree` and `lxml` — every shape is backed by a real `oxml` element.
Prefer the high-level API (`shape.text_frame`, `shape.fill.solid()`,
`shape.left`/`.top`/`.width`/`.height`) whenever it covers what you need; reach for
raw XML only when it doesn't.

## Thumbnails & Analysis
There is no slide-rendering engine in the sandbox (no LibreOffice/PowerPoint) — do
NOT claim a pixel-perfect visual render. Instead:
- **Structural analysis** (the reliable part): iterate `slide.shapes` for shape
  counts, `shape.shape_type`, positions/sizes, `shape.has_chart`/`has_table`, and
  `shape.text_frame.text` for all text content — this is exact, not approximate.
- **Schematic thumbnail** (a wireframe, not a real render): draw each shape's
  bounding box and a text snippet with `Pillow` (`ImageDraw.rectangle` +
  `ImageDraw.text`, scaled from EMU to pixels). Label the output clearly as a
  schematic/wireframe overview when you show it to the user — it will not match
  the deck's actual fonts, colors, or chart rendering.
"""

_XLSX_CONTENT = """\
Build, analyze, and visualize Excel workbooks with `openpyxl` and `pandas` (already
installed in the sandbox) via the `exec` tool. Write outputs to /workspace.

## Spreadsheet Creation
Write two small local helper functions at the top of the script — `build_table` and
`build_chart` — parameterized by a `theme` (a dict of header/accent colors) so every
sheet in the workbook looks consistent without repeating style code:

```python
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import BarChart, Reference

THEMES = {
    "corporate": {"header_fill": "1F4E79", "header_font": "FFFFFF", "accent": "2E75B6"},
    "warm":      {"header_fill": "C0504D", "header_font": "FFFFFF", "accent": "E8A33D"},
}

def build_table(ws, rows, start_row=1, theme="corporate"):
    t = THEMES[theme]
    header_fill = PatternFill("solid", fgColor=t["header_fill"])
    header_font = Font(bold=True, color=t["header_font"])
    for c, header in enumerate(rows[0], start=1):
        cell = ws.cell(row=start_row, column=c, value=header)
        cell.fill, cell.font = header_fill, header_font
    for r, row in enumerate(rows[1:], start=start_row + 1):
        for c, value in enumerate(row, start=1):
            ws.cell(row=r, column=c, value=value)
    return start_row + len(rows)  # next free row

def build_chart(ws, data_range, cats_range, anchor, title, theme="corporate", kind="bar"):
    chart = BarChart() if kind == "bar" else BarChart()  # swap for LineChart/AreaChart as needed
    chart.title = title
    chart.add_data(data_range, titles_from_data=True)
    chart.set_categories(cats_range)
    ws.add_chart(chart, anchor)

wb = Workbook()
ws = wb.active
next_row = build_table(ws, [["Supplier", "Q1", "Q2", "Q3"], ["Acme", 100, 120, 115]])
# Formula cells: write the formula STRING, not a computed value — Excel/LibreOffice
# calculates it on open. openpyxl never evaluates formulas itself.
ws["E2"] = "=SUM(B2:D2)"
build_chart(ws, Reference(ws, min_col=2, min_row=1, max_col=4, max_row=2),
            Reference(ws, min_col=1, min_row=2, max_row=2), "G2", "Quarterly spend")
wb.save("/workspace/suppliers.xlsx")
```
For **line/area/combo charts** swap `BarChart` for `openpyxl.chart.LineChart` /
`AreaChart`, or add a second series with `chart2.y_axis.axId = chart.y_axis.axId` and
`chart += chart2` for a combo. Configure axis titles via `chart.x_axis.title` /
`chart.y_axis.title`, and legend position via `chart.legend.position`.

## Data Analysis
Use `pandas` for anything beyond simple rows: `pd.read_excel`/`read_csv` to load,
`.groupby()`/`.describe()`/`.pivot_table()` for statistics, then either
`df.to_excel(path, sheet_name=...)` for a quick dump, or `df.values.tolist()` fed
into `build_table` above when you need the openpyxl styling/formulas on top.

## Financial Models
For DCF / LBO / three-statement models, the point is that every number the reader
might question should be a **formula referencing another cell**, not a hardcoded
literal — so changing one assumption recalculates the whole model:
- Put assumptions (discount rate, growth rate, tax rate, leverage ratio) in one
  labelled "Assumptions" block/sheet, each in its own cell.
- Every downstream calculation references those cells (`=B5*(1+$B$2)`), never
  repeats the number.
- **Source annotations**: add a cell comment (`ws["B2"].comment = Comment("Source: ...", "Claw")`)
  or a footnote row citing where each assumption came from.
- **Scenario analysis**: put 2-3 named scenarios (Base/Upside/Downside) as columns
  of assumption values on a separate sheet, and have the model sheet's assumption
  cells reference one scenario column at a time (or use Excel Data Tables /
  `CHOOSE()` if the target is manual toggling after you hand off the file).
"""

_PDF_CONTENT = """\
Extract, create, merge, and fill PDFs with `pypdf` and `PyMuPDF` (`fitz`) — both
pre-installed in the sandbox — via the `exec` tool. Prefer `pypdf` for
merge/split/form-fill and `PyMuPDF` for extraction-with-layout and form-field
creation; both read/write the same PDF files so mix freely.

## Text & Table Extraction
`PyMuPDF` preserves layout far better than plain text scraping:

```python
import fitz  # PyMuPDF

doc = fitz.open("/workspace/contract.pdf")
for page in doc:
    text = page.get_text("text")       # reading-order plain text
    words = page.get_text("words")     # (x0, y0, x1, y1, word, block, line, word_no)
    tables = page.find_tables()        # page.find_tables().tables -> structured cells
    for table in tables.tables:
        rows = table.extract()         # list[list[str]] — feed straight into pandas/csv
```
Use `words`/block positions when you need to preserve columns or find text near a
specific location (e.g. "the clause after this heading"), not just concatenated text.

## PDF Creation
For a **document from scratch** (reports, letters), `reportlab` is more direct for
flowing text/tables (`SimpleDocTemplate` + `Paragraph`/`Table` flowables). For
**precise, coordinate-based placement** (a branded one-pager, an overlay on an
existing template), build pages with `fitz`:

```python
doc = fitz.open()
page = doc.new_page(width=595, height=842)  # A4 in points
page.insert_text((72, 72), "Purchase Order #4471", fontsize=18, fontname="helv")
page.draw_rect(fitz.Rect(72, 100, 523, 101), color=(0.2, 0.3, 0.5), fill=(0.2, 0.3, 0.5))
doc.save("/workspace/po.pdf")
```
Both libraries support multi-page docs (call the page-adding API in a loop) and
embedding chart images (render the chart with matplotlib/openpyxl to a PNG first,
then `page.insert_image(rect, filename=...)` or reportlab's `Image` flowable).

## Merge & Split
```python
from pypdf import PdfWriter, PdfReader

writer = PdfWriter()
writer.append("/workspace/cover.pdf")
writer.append("/workspace/contract.pdf")
writer.write("/workspace/combined.pdf")

reader = PdfReader("/workspace/combined.pdf")
writer = PdfWriter()
writer.add_page(reader.pages[2].rotate(90))       # rotate a specific page
writer.add_page(reader.pages[0])                  # reorder — page 0 goes last
writer.write("/workspace/reordered.pdf")
```

## Form Filling
Create a fillable field with `PyMuPDF`, then fill it with `pypdf` (the two libraries
round-trip AcroForm fields cleanly):

```python
doc = fitz.open("/workspace/template.pdf")
widget = fitz.Widget()
widget.field_name, widget.field_type = "supplier_name", fitz.PDF_WIDGET_TYPE_TEXT
widget.rect = fitz.Rect(150, 700, 400, 720)
doc[0].add_widget(widget)
doc.save("/workspace/form.pdf")

reader = PdfReader("/workspace/form.pdf")
writer = PdfWriter()
writer.append(reader)
writer.update_page_form_field_values(writer.pages[0], {"supplier_name": "Acme Co."})
writer.write("/workspace/filled.pdf")
```
Before filling an unfamiliar PDF, first detect its existing fields with
`reader.get_fields()` (pypdf) to confirm exact field names, then validate required
fields got a non-empty value before handing the filled PDF back to the user.
"""

_DOCX_CONTENT = """\
Create and edit Word documents with `python-docx` (already installed in the sandbox)
via the `exec` tool. Write outputs to /workspace.

## Document Creation
```python
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

doc = Document()
doc.styles["Normal"].font.name, doc.styles["Normal"].font.size = "Calibri", Pt(11)

h = doc.add_heading("Request for Quotation", level=1)
p = doc.add_paragraph("Dear Supplier,")
table = doc.add_table(rows=1, cols=3)
table.style = "Light Grid Accent 1"
table.rows[0].cells[0].text = "Item"

section = doc.sections[0]
section.header.paragraphs[0].text = "Acme Procurement — Confidential"
doc.save("/workspace/rfq.docx")
```
For professional typesetting: set styles at the document level (above) rather than
per-run, use built-in heading styles (`add_heading`) so a generated table of contents
picks them up, and `section.page_width`/`page_height`/margins for layout control.

## Document Editing
When editing an EXISTING document, preserve its formatting: edit `run.text` on the
specific run that needs to change rather than replacing `paragraph.text` (which
destroys any bold/italic/color split across runs in that paragraph). For structural
edits python-docx's API doesn't expose (custom OOXML attributes), drop to
`paragraph._p` (the underlying `lxml` element) directly — see Tracked Changes below
for the pattern.

## Tracked Changes
`python-docx` has **no native tracked-changes API** — Word's revision marks
(`w:ins`/`w:del`) must be built as raw OOXML and inserted via the low-level
`oxml`/`lxml` layer:

```python
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import datetime

def insert_tracked_insertion(paragraph, text, author="Claw"):
    ins = OxmlElement("w:ins")
    ins.set(qn("w:id"), "1")
    ins.set(qn("w:author"), author)
    ins.set(qn("w:date"), datetime.datetime.now(datetime.timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"))
    run = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.set(qn("xml:space"), "preserve")
    t.text = text
    run.append(t)
    ins.append(run)
    paragraph._p.append(ins)

def mark_tracked_deletion(run, author="Claw"):
    # Wrap the run's text as w:delText inside a w:del, replacing the plain w:r.
    del_el = OxmlElement("w:del")
    del_el.set(qn("w:id"), "2")
    del_el.set(qn("w:author"), author)
    del_el.set(qn("w:date"), datetime.datetime.now(datetime.timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"))
    new_run = OxmlElement("w:r")
    del_text = OxmlElement("w:delText")
    del_text.set(qn("xml:space"), "preserve")
    del_text.text = run.text
    new_run.append(del_text)
    del_el.append(new_run)
    run._element.getparent().replace(run._element, del_el)
```
This is exactly what Word renders as "tracked changes" in Review mode — suitable for
legal/academic/business review workflows where the user needs to see (and accept or
reject) each change rather than a silently-edited document. Review comments follow
the same low-level pattern via the `w:commentRangeStart`/`w:comment` parts.

## Content Extraction
```python
doc = Document("/workspace/contract.docx")
text = "\\n".join(p.text for p in doc.paragraphs)
tables = [[[cell.text for cell in row.cells] for row in t.rows] for t in doc.tables]
metadata = doc.core_properties  # .author, .title, .created, ...

# Embedded images live in the document part's relationships:
for rel in doc.part.rels.values():
    if "image" in rel.reltype:
        image_bytes = rel.target_part.blob
```
Comments require the same low-level `oxml` route as tracked changes (there's no
`doc.comments` on python-docx) — read the `word/comments.xml` part via
`doc.part.package.part_related_by(...)` if the document has one. Support multiple
output formats (plain text, markdown, JSON) by building a plain dict/list from the
above and letting the caller decide the serialization.
"""

_FINANCIAL_STATEMENT_CONTENT = r"""
Analyze a company's income statement, balance sheet, and cash flow statement:
compute year-over-year (YoY) and quarter-over-quarter (QoQ) changes, financial
ratios, and run 10 built-in anomaly-detection rules (AR surge, cash flow
divergence, inventory buildup, margin shifts, excessive leverage/goodwill,
weak liquidity, AP anomalies), then produce a readable report.

## Output Language
Write the final report in Thai (ภาษาไทย) by default. Only use a different
language if the user has explicitly asked for one (in this request or earlier
in the conversation). Numbers, financial terms, and the underlying script's
CLI flags/JSON keys stay as-is — translate the surrounding narrative and
section headings, not the data itself.

## Setup (once per workspace)
This skill ships with a complete, dependency-free (stdlib-only) analysis
script rather than reimplementing ~700 lines of ratio/anomaly logic from
scratch on every use — more reliable and far cheaper in tokens. Before first
use in a workspace, check whether it already exists:

```bash
test -f /workspace/.claw_skills/financial-statement-analyzer/analyze_financials.py && echo present
```

If missing, create the directory and use `write_file` to save the script
below verbatim to
`/workspace/.claw_skills/financial-statement-analyzer/analyze_financials.py`.
Do this once per workspace; every later analysis just calls the existing file
via `exec`.

```python
#!/usr/bin/env python3
'''
Financial Statement Analyzer
Computes YoY/QoQ changes, financial ratios, and anomaly detection.
'''

import argparse
import json
import re
import sys
from collections import OrderedDict


SAMPLE_DATA = {
    "company": "Acme Corp",
    "currency": "USD",
    "unit": "thousands",
    "periods": [
        "2023Q1", "2023Q2", "2023Q3", "2023Q4",
        "2024Q1", "2024Q2", "2024Q3", "2024Q4"
    ],
    "income_statement": {
        "revenue":         [5000, 5200, 4800, 6000, 5500, 5800, 5100, 6500],
        "cost_of_revenue": [3000, 3100, 2900, 3500, 3400, 3600, 3200, 4100],
        "operating_income": [800,  850,  750, 1000,  780,  820,  700,  900],
        "net_income":       [600,  650,  560,  780,  580,  620,  520,  680]
    },
    "balance_sheet": {
        "accounts_receivable":    [2000, 2100, 2200, 2300, 2800, 3200, 3600, 4200],
        "inventory":              [1000, 1050, 1100, 1200, 1100, 1150, 1200, 1300],
        "total_current_assets":   [5000, 5200, 5400, 5800, 6000, 6500, 7000, 7500],
        "goodwill":               [ 500,  500,  500,  500,  500,  500,  500,  500],
        "total_assets":           [15000,15500,16000,16500,17000,17500,18000,18500],
        "accounts_payable":       [1500, 1600, 1550, 1700, 1650, 1750, 1700, 1800],
        "total_current_liabilities":[4000,4200,4100,4500,4300,4600,4500,4900],
        "total_liabilities":      [8000, 8200, 8400, 8600, 8800, 9000, 9200, 9500],
        "total_equity":           [7000, 7300, 7600, 7900, 8200, 8500, 8800, 9000]
    },
    "cash_flow": {
        "operating_cash_flow": [ 700,  750,  620,  850,  300,  280,  250,  200],
        "investing_cash_flow": [-200, -180, -250, -300, -400, -350, -300, -280],
        "financing_cash_flow": [-100,  -50,  -80, -120,  200,  150,  100,   50],
        "capex":               [ 180,  160,  230,  280,  380,  330,  280,  260]
    }
}

DEFAULT_THRESHOLDS = {
    "ar_threshold": 0.20,
    "inv_threshold": 0.15,
    "ocf_ratio": 0.50,
    "margin_threshold": 0.05,
    "net_margin_threshold": 0.03,
    "debt_ceiling": 0.70,
    "current_floor": 1.00,
    "goodwill_ceiling": 0.30,
    "neg_ocf_periods": 2,
    "ap_threshold": 0.20,
}


def safe_div(a, b):
    if b is None or a is None:
        return None
    if b == 0:
        return None
    return a / b


def pct_change(new, old):
    if old is None or new is None:
        return None
    if old == 0:
        if new == 0:
            return 0.0
        return None
    return (new - old) / abs(old)


def get_values(data, section, key):
    sec = data.get(section, {})
    return sec.get(key)


def detect_period_type(periods):
    for p in periods:
        if re.search(r"[Qq]\d", str(p)):
            return "quarterly"
    return "annual"


def parse_period(p):
    p = str(p).strip()
    m = re.match(r"(\d{4})[Qq](\d)", p)
    if m:
        return int(m.group(1)), int(m.group(2))
    m2 = re.match(r"(\d{4})", p)
    if m2:
        return int(m2.group(1)), 0
    return None, None


def build_yoy_map(periods):
    idx = {}
    for i, p in enumerate(periods):
        year, q = parse_period(p)
        if year is not None:
            idx[(year, q)] = i

    yoy_pairs = []
    for i, p in enumerate(periods):
        year, q = parse_period(p)
        if year is None:
            yoy_pairs.append(None)
            continue
        prev_key = (year - 1, q)
        if prev_key in idx:
            yoy_pairs.append(idx[prev_key])
        else:
            yoy_pairs.append(None)
    return yoy_pairs


def compute_changes(values, periods, yoy_map):
    n = len(values)
    yoy = []
    qoq = []
    for i in range(n):
        if yoy_map[i] is not None:
            yoy.append(pct_change(values[i], values[yoy_map[i]]))
        else:
            yoy.append(None)

        if i > 0:
            qoq.append(pct_change(values[i], values[i - 1]))
        else:
            qoq.append(None)
    return yoy, qoq


def compute_all_changes(data, periods, yoy_map):
    result = {}
    for section in ("income_statement", "balance_sheet", "cash_flow"):
        sec = data.get(section, {})
        for key, vals in sec.items():
            if not isinstance(vals, list):
                continue
            yoy, qoq = compute_changes(vals, periods, yoy_map)
            result[key] = {"yoy": yoy, "qoq": qoq}
    return result


def compute_ratios(data, periods):
    n = len(periods)
    ratios = []
    for i in range(n):
        r = OrderedDict()
        rev = _val(data, "income_statement", "revenue", i)
        cogs = _val(data, "income_statement", "cost_of_revenue", i)
        oi = _val(data, "income_statement", "operating_income", i)
        ni = _val(data, "income_statement", "net_income", i)
        ar = _val(data, "balance_sheet", "accounts_receivable", i)
        inv = _val(data, "balance_sheet", "inventory", i)
        tca = _val(data, "balance_sheet", "total_current_assets", i)
        gw = _val(data, "balance_sheet", "goodwill", i)
        ta = _val(data, "balance_sheet", "total_assets", i)
        ap = _val(data, "balance_sheet", "accounts_payable", i)
        tcl = _val(data, "balance_sheet", "total_current_liabilities", i)
        tl = _val(data, "balance_sheet", "total_liabilities", i)
        te = _val(data, "balance_sheet", "total_equity", i)
        ocf = _val(data, "cash_flow", "operating_cash_flow", i)
        capex = _val(data, "cash_flow", "capex", i)

        gp = (rev - cogs) if (rev is not None and cogs is not None) else None
        r["gross_margin"] = safe_div(gp, rev)
        r["operating_margin"] = safe_div(oi, rev)
        r["net_margin"] = safe_div(ni, rev)
        r["debt_ratio"] = safe_div(tl, ta)
        r["current_ratio"] = safe_div(tca, tcl)
        r["goodwill_ratio"] = safe_div(gw, ta)
        r["ar_to_revenue"] = safe_div(ar, rev)
        r["inventory_to_revenue"] = safe_div(inv, rev)
        r["ocf_to_ni"] = safe_div(ocf, ni) if ni and ni != 0 else None
        r["roe"] = safe_div(ni, te)

        fcf = None
        if ocf is not None and capex is not None:
            fcf = ocf - capex
        r["free_cash_flow"] = fcf

        dso = None
        if ar is not None and rev is not None and rev != 0:
            dso = (ar / rev) * 90 if detect_period_type(periods) == "quarterly" else (ar / rev) * 365
        r["dso"] = dso

        dio = None
        if inv is not None and cogs is not None and cogs != 0:
            dio = (inv / cogs) * 90 if detect_period_type(periods) == "quarterly" else (inv / cogs) * 365
        r["dio"] = dio

        dpo = None
        if ap is not None and cogs is not None and cogs != 0:
            dpo = (ap / cogs) * 90 if detect_period_type(periods) == "quarterly" else (ap / cogs) * 365
        r["dpo"] = dpo

        ratios.append(r)
    return ratios


def _val(data, section, key, idx):
    vals = get_values(data, section, key)
    if vals is None or idx >= len(vals):
        return None
    v = vals[idx]
    return v if v is not None else None


def detect_anomalies(data, periods, changes, ratios, thresholds):
    alerts = []
    n = len(periods)

    rev_yoy = changes.get("revenue", {}).get("yoy", [])
    ar_yoy = changes.get("accounts_receivable", {}).get("yoy", [])
    if rev_yoy and ar_yoy:
        for i in range(n):
            if ar_yoy[i] is not None and rev_yoy[i] is not None:
                gap = ar_yoy[i] - rev_yoy[i]
                if gap > thresholds["ar_threshold"]:
                    alerts.append({
                        "rule": "AR Surge",
                        "period": periods[i],
                        "severity": "high" if gap > thresholds["ar_threshold"] * 2 else "medium",
                        "detail": (
                            f"AR growth {ar_yoy[i]*100:.1f}% far outpaces revenue growth "
                            f"{rev_yoy[i]*100:.1f}%, gap {gap*100:.1f}pp (threshold "
                            f"{thresholds['ar_threshold']*100:.0f}pp)"
                        ),
                        "implication": "Possible aggressive revenue recognition or collection difficulties"
                    })

    for i in range(n):
        ni = _val(data, "income_statement", "net_income", i)
        ocf = _val(data, "cash_flow", "operating_cash_flow", i)
        if ni is not None and ocf is not None and ni != 0:
            ratio = ocf / ni if ni > 0 else None
            if ni > 0 and ocf < 0:
                alerts.append({
                    "rule": "Cash Flow Divergence",
                    "period": periods[i],
                    "severity": "high",
                    "detail": f"Net income {ni:,.0f} is positive but operating cash flow {ocf:,.0f} is negative",
                    "implication": "Low earnings quality; profits may contain significant accruals"
                })
            elif ratio is not None and ratio < thresholds["ocf_ratio"]:
                alerts.append({
                    "rule": "Cash Flow Divergence",
                    "period": periods[i],
                    "severity": "medium",
                    "detail": (
                        f"OCF/Net Income = {ratio:.2f}, below threshold {thresholds['ocf_ratio']:.2f}"
                    ),
                    "implication": "Profit-to-cash conversion efficiency is low"
                })

    inv_yoy = changes.get("inventory", {}).get("yoy", [])
    if rev_yoy and inv_yoy:
        for i in range(n):
            if inv_yoy[i] is not None and rev_yoy[i] is not None:
                gap = inv_yoy[i] - rev_yoy[i]
                if gap > thresholds["inv_threshold"]:
                    alerts.append({
                        "rule": "Inventory Buildup",
                        "period": periods[i],
                        "severity": "medium",
                        "detail": (
                            f"Inventory growth {inv_yoy[i]*100:.1f}% exceeds revenue growth "
                            f"{rev_yoy[i]*100:.1f}%, gap {gap*100:.1f}pp"
                        ),
                        "implication": "Potential product obsolescence or write-down risk"
                    })

    for i in range(n):
        if i < 1:
            continue
        gm_now = ratios[i].get("gross_margin")
        gm_prev = ratios[i - 1].get("gross_margin")
        if gm_now is not None and gm_prev is not None:
            delta = gm_now - gm_prev
            if abs(delta) > thresholds["margin_threshold"]:
                direction = "up" if delta > 0 else "down"
                alerts.append({
                    "rule": "Gross Margin Shift",
                    "period": periods[i],
                    "severity": "medium",
                    "detail": (
                        f"Gross margin {direction} {abs(delta)*100:.1f}pp QoQ "
                        f"({gm_prev*100:.1f}% -> {gm_now*100:.1f}%)"
                    ),
                    "implication": "Significant change in pricing power or cost structure"
                })

    for i in range(n):
        if i < 1:
            continue
        nm_now = ratios[i].get("net_margin")
        nm_prev = ratios[i - 1].get("net_margin")
        if nm_now is not None and nm_prev is not None:
            delta = nm_now - nm_prev
            if abs(delta) > thresholds["net_margin_threshold"]:
                direction = "up" if delta > 0 else "down"
                alerts.append({
                    "rule": "Net Margin Shift",
                    "period": periods[i],
                    "severity": "medium",
                    "detail": (
                        f"Net margin {direction} {abs(delta)*100:.1f}pp QoQ "
                        f"({nm_prev*100:.1f}% -> {nm_now*100:.1f}%)"
                    ),
                    "implication": "Abnormal expense control or non-recurring items"
                })

    streak = 0
    for i in range(n):
        ocf = _val(data, "cash_flow", "operating_cash_flow", i)
        if ocf is not None and ocf < 0:
            streak += 1
            if streak >= thresholds["neg_ocf_periods"]:
                alerts.append({
                    "rule": "Persistent Negative OCF",
                    "period": periods[i],
                    "severity": "high" if streak >= 3 else "medium",
                    "detail": f"Operating cash flow has been negative for {streak} consecutive periods",
                    "implication": "Insufficient organic cash generation; reliant on external financing"
                })
        else:
            streak = 0

    for i in range(n):
        gw_ratio = ratios[i].get("goodwill_ratio")
        if gw_ratio is not None and gw_ratio > thresholds["goodwill_ceiling"]:
            alerts.append({
                "rule": "Excessive Goodwill",
                "period": periods[i],
                "severity": "medium",
                "detail": f"Goodwill/Total Assets = {gw_ratio*100:.1f}% (ceiling {thresholds['goodwill_ceiling']*100:.0f}%)",
                "implication": "Impairment risk if acquired entities underperform"
            })

    for i in range(n):
        dr = ratios[i].get("debt_ratio")
        if dr is not None and dr > thresholds["debt_ceiling"]:
            alerts.append({
                "rule": "High Leverage",
                "period": periods[i],
                "severity": "medium",
                "detail": f"Debt-to-asset ratio {dr*100:.1f}% (ceiling {thresholds['debt_ceiling']*100:.0f}%)",
                "implication": "Elevated financial leverage; higher debt repayment pressure"
            })

    for i in range(n):
        cr = ratios[i].get("current_ratio")
        if cr is not None and cr < thresholds["current_floor"]:
            alerts.append({
                "rule": "Low Current Ratio",
                "period": periods[i],
                "severity": "medium" if cr > 0.7 else "high",
                "detail": f"Current ratio {cr:.2f} (floor {thresholds['current_floor']:.2f})",
                "implication": "Weak short-term liquidity; potential liquidity risk"
            })

    cogs_yoy = changes.get("cost_of_revenue", {}).get("yoy", [])
    ap_yoy = changes.get("accounts_payable", {}).get("yoy", [])
    if cogs_yoy and ap_yoy:
        for i in range(n):
            if ap_yoy[i] is not None and cogs_yoy[i] is not None:
                gap = abs(ap_yoy[i] - cogs_yoy[i])
                if gap > thresholds["ap_threshold"]:
                    if ap_yoy[i] > cogs_yoy[i]:
                        msg = "AP growing much faster than COGS — may be stretching payment terms to ease cash pressure"
                    else:
                        msg = "AP growing much slower than COGS — suppliers may be demanding shorter payment terms"
                    alerts.append({
                        "rule": "AP Anomaly",
                        "period": periods[i],
                        "severity": "low",
                        "detail": (
                            f"AP growth {ap_yoy[i]*100:.1f}% vs COGS growth {cogs_yoy[i]*100:.1f}%, "
                            f"gap {gap*100:.1f}pp"
                        ),
                        "implication": msg
                    })

    return alerts


def fmt_pct(v):
    if v is None:
        return "N/A"
    return f"{v * 100:+.1f}%"


def fmt_ratio(v, decimals=2):
    if v is None:
        return "N/A"
    return f"{v:.{decimals}f}"


def fmt_num(v):
    if v is None:
        return "N/A"
    return f"{v:,.0f}"


def format_text_report(data, periods, changes, ratios, alerts):
    lines = []
    company = data.get("company", "Unknown Company")
    unit = data.get("unit", "")
    currency = data.get("currency", "")
    unit_label = f"({currency} {unit})" if currency or unit else ""

    lines.append(f"{'='*60}")
    lines.append(f"  Financial Statement Analysis Report — {company}")
    lines.append(f"  Period: {periods[0]} ~ {periods[-1]} {unit_label}")
    lines.append(f"{'='*60}")
    lines.append("")

    lines.append("[1. Key Metrics at a Glance (Latest Period)]")
    lines.append("")
    latest = ratios[-1]
    prev = ratios[-2] if len(ratios) > 1 else {}
    metrics = [
        ("Gross Margin", "gross_margin"),
        ("Operating Margin", "operating_margin"),
        ("Net Margin", "net_margin"),
        ("Debt Ratio", "debt_ratio"),
        ("Current Ratio", "current_ratio"),
        ("Goodwill Ratio", "goodwill_ratio"),
        ("AR/Revenue", "ar_to_revenue"),
        ("OCF/Net Income", "ocf_to_ni"),
        ("ROE", "roe"),
    ]
    lines.append(f"  {'Metric':<18} {'Current':>10} {'Prior':>10} {'Change':>10}")
    lines.append(f"  {'-'*18} {'-'*10} {'-'*10} {'-'*10}")
    for label, key in metrics:
        curr_v = latest.get(key)
        prev_v = prev.get(key) if prev else None
        delta = None
        if curr_v is not None and prev_v is not None:
            delta = curr_v - prev_v
        c_str = fmt_ratio(curr_v) if key == "current_ratio" else (fmt_pct(curr_v).replace("+", "") if curr_v is not None else "N/A")
        p_str = fmt_ratio(prev_v) if key == "current_ratio" else (fmt_pct(prev_v).replace("+", "") if prev_v is not None else "N/A")
        d_str = fmt_pct(delta) if delta is not None else "N/A"
        if key == "current_ratio":
            c_str = fmt_ratio(curr_v)
            p_str = fmt_ratio(prev_v)
            d_str = fmt_ratio(delta) if delta is not None else "N/A"
        lines.append(f"  {label:<18} {c_str:>10} {p_str:>10} {d_str:>10}")

    lines.append("")
    lines.append(f"  Days Sales Outstanding (DSO): {fmt_ratio(latest.get('dso'), 0)} days")
    lines.append(f"  Days Inventory Outstanding (DIO): {fmt_ratio(latest.get('dio'), 0)} days")
    lines.append(f"  Days Payable Outstanding (DPO): {fmt_ratio(latest.get('dpo'), 0)} days")
    fcf = latest.get("free_cash_flow")
    lines.append(f"  Free Cash Flow (FCF):   {fmt_num(fcf)}")
    lines.append("")

    lines.append("[2. Year-over-Year Changes (YoY)]")
    lines.append("")
    yoy_keys = [
        ("Revenue", "revenue"), ("Cost of Revenue", "cost_of_revenue"),
        ("Operating Income", "operating_income"), ("Net Income", "net_income"),
        ("Accounts Receivable", "accounts_receivable"), ("Inventory", "inventory"),
        ("Total Assets", "total_assets"), ("Total Liabilities", "total_liabilities"),
        ("Operating Cash Flow", "operating_cash_flow"),
    ]
    header = f"  {'Line Item':<20}"
    for p in periods:
        header += f" {p:>9}"
    lines.append(header)
    lines.append(f"  {'-'*20}" + f" {'-'*9}" * len(periods))
    for label, key in yoy_keys:
        ch = changes.get(key, {}).get("yoy", [])
        if not ch:
            continue
        row = f"  {label:<20}"
        for v in ch:
            row += f" {fmt_pct(v):>9}"
        lines.append(row)
    lines.append("")

    lines.append("[3. Quarter-over-Quarter Changes (QoQ)]")
    lines.append("")
    header = f"  {'Line Item':<20}"
    for p in periods:
        header += f" {p:>9}"
    lines.append(header)
    lines.append(f"  {'-'*20}" + f" {'-'*9}" * len(periods))
    for label, key in yoy_keys:
        ch = changes.get(key, {}).get("qoq", [])
        if not ch:
            continue
        row = f"  {label:<20}"
        for v in ch:
            row += f" {fmt_pct(v):>9}"
        lines.append(row)
    lines.append("")

    lines.append("[4. Anomaly Detection Results]")
    lines.append("")
    if not alerts:
        lines.append("  No anomalies detected.")
    else:
        severity_order = {"high": 0, "medium": 1, "low": 2}
        sorted_alerts = sorted(alerts, key=lambda a: (severity_order.get(a["severity"], 9), a["period"]))
        severity_labels = {"high": "High", "medium": "Medium", "low": "Low"}
        for idx, a in enumerate(sorted_alerts, 1):
            sev = severity_labels.get(a["severity"], a["severity"])
            lines.append(f"  [{sev} risk] {a['rule']} | {a['period']}")
            lines.append(f"    Detail: {a['detail']}")
            lines.append(f"    Implication: {a['implication']}")
            lines.append("")
    lines.append("")

    lines.append("[5. Cross-Statement Analysis Highlights]")
    lines.append("")
    rev_vals = get_values(data, "income_statement", "revenue")
    ar_vals = get_values(data, "balance_sheet", "accounts_receivable")
    ni_vals = get_values(data, "income_statement", "net_income")
    ocf_vals = get_values(data, "cash_flow", "operating_cash_flow")

    if rev_vals and ar_vals and len(rev_vals) >= 2:
        rev_growth = pct_change(rev_vals[-1], rev_vals[-2])
        ar_growth = pct_change(ar_vals[-1], ar_vals[-2])
        if rev_growth is not None and ar_growth is not None:
            if ar_growth > rev_growth + 0.1:
                lines.append("  - Income Statement -> Balance Sheet: AR growth notably outpaces revenue growth; revenue growth quality warrants attention")
            else:
                lines.append("  - Income Statement -> Balance Sheet: AR and revenue growth are broadly in line; revenue growth quality looks reasonable")

    if ni_vals and ocf_vals:
        total_ni = sum(v for v in ni_vals if v is not None)
        total_ocf = sum(v for v in ocf_vals if v is not None)
        if total_ni > 0:
            overall_ratio = total_ocf / total_ni
            if overall_ratio < 0.6:
                lines.append(f"  - Income Statement -> Cash Flow Statement: cumulative OCF/Net Income = {overall_ratio:.2f}, cash content of earnings is low")
            else:
                lines.append(f"  - Income Statement -> Cash Flow Statement: cumulative OCF/Net Income = {overall_ratio:.2f}, cash content of earnings is reasonable")

    inv_cf = get_values(data, "cash_flow", "investing_cash_flow")
    capex_vals = get_values(data, "cash_flow", "capex")
    if inv_cf and capex_vals:
        latest_inv = inv_cf[-1] if inv_cf[-1] is not None else 0
        latest_capex = capex_vals[-1] if capex_vals[-1] is not None else 0
        if latest_inv is not None and latest_capex is not None:
            lines.append(
                f"  - Balance Sheet -> Cash Flow Statement: investing cash flow {fmt_num(latest_inv)}, "
                f"capex {fmt_num(latest_capex)}"
            )

    lines.append("")
    lines.append(f"{'='*60}")
    lines.append("  Report complete")
    lines.append(f"{'='*60}")
    return "\n".join(lines)


def build_json_report(data, periods, changes, ratios, alerts):
    return OrderedDict([
        ("company", data.get("company", "")),
        ("report_range", f"{periods[0]} ~ {periods[-1]}"),
        ("periods", periods),
        ("ratios", [dict(r) for r in ratios]),
        ("changes", {k: dict(v) for k, v in changes.items()}),
        ("anomalies", alerts),
        ("summary", {
            "total_anomalies": len(alerts),
            "high_severity": sum(1 for a in alerts if a["severity"] == "high"),
            "medium_severity": sum(1 for a in alerts if a["severity"] == "medium"),
            "low_severity": sum(1 for a in alerts if a["severity"] == "low"),
        })
    ])


def validate_data(data):
    errors = []
    periods = data.get("periods")
    if not periods or not isinstance(periods, list):
        errors.append("missing 'periods' field or wrong format")
        return errors

    n = len(periods)
    for section in ("income_statement", "balance_sheet", "cash_flow"):
        sec = data.get(section, {})
        for key, vals in sec.items():
            if isinstance(vals, list) and len(vals) != n:
                errors.append(
                    f"{section}.{key} has length {len(vals)}, expected {n} (must match periods)"
                )
    return errors


def main():
    parser = argparse.ArgumentParser(
        description="Financial statement analyzer — YoY/QoQ trends + anomaly detection",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Examples:\n"
               "  %(prog)s data.json\n"
               "  %(prog)s data.json --json\n"
               "  %(prog)s data.json -o report.json\n"
               "  %(prog)s --sample > sample.json\n"
    )
    parser.add_argument("input", nargs="?", help="path to the input JSON file")
    parser.add_argument("-j", "--json", action="store_true", help="output in JSON format")
    parser.add_argument("-o", "--output", help="output file path")
    parser.add_argument("-s", "--sample", action="store_true", help="print sample data to stdout")

    parser.add_argument("--ar-threshold", type=float, default=DEFAULT_THRESHOLDS["ar_threshold"],
                        help=f"AR anomaly threshold (default {DEFAULT_THRESHOLDS['ar_threshold']})")
    parser.add_argument("--inv-threshold", type=float, default=DEFAULT_THRESHOLDS["inv_threshold"],
                        help=f"inventory anomaly threshold (default {DEFAULT_THRESHOLDS['inv_threshold']})")
    parser.add_argument("--ocf-ratio", type=float, default=DEFAULT_THRESHOLDS["ocf_ratio"],
                        help=f"cash flow / profit divergence threshold (default {DEFAULT_THRESHOLDS['ocf_ratio']})")
    parser.add_argument("--margin-threshold", type=float, default=DEFAULT_THRESHOLDS["margin_threshold"],
                        help=f"gross margin shift threshold (default {DEFAULT_THRESHOLDS['margin_threshold']})")
    parser.add_argument("--debt-ceiling", type=float, default=DEFAULT_THRESHOLDS["debt_ceiling"],
                        help=f"debt-to-asset ratio warning level (default {DEFAULT_THRESHOLDS['debt_ceiling']})")
    parser.add_argument("--current-floor", type=float, default=DEFAULT_THRESHOLDS["current_floor"],
                        help=f"current ratio warning level (default {DEFAULT_THRESHOLDS['current_floor']})")
    parser.add_argument("--goodwill-ceiling", type=float, default=DEFAULT_THRESHOLDS["goodwill_ceiling"],
                        help=f"goodwill-to-asset ratio warning level (default {DEFAULT_THRESHOLDS['goodwill_ceiling']})")

    args = parser.parse_args()

    if args.sample:
        json.dump(SAMPLE_DATA, sys.stdout, ensure_ascii=False, indent=2)
        sys.stdout.write("\n")
        return

    if not args.input:
        parser.error("provide an input JSON file path, or use --sample to generate sample data")

    try:
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
    except FileNotFoundError:
        print(f"Error: file '{args.input}' not found", file=sys.stderr)
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: failed to parse JSON — {e}", file=sys.stderr)
        sys.exit(1)

    errors = validate_data(data)
    if errors:
        print("Data validation failed:", file=sys.stderr)
        for err in errors:
            print(f"  - {err}", file=sys.stderr)
        sys.exit(1)

    periods = data["periods"]
    yoy_map = build_yoy_map(periods)

    thresholds = {
        "ar_threshold": args.ar_threshold,
        "inv_threshold": args.inv_threshold,
        "ocf_ratio": args.ocf_ratio,
        "margin_threshold": args.margin_threshold,
        "net_margin_threshold": DEFAULT_THRESHOLDS["net_margin_threshold"],
        "debt_ceiling": args.debt_ceiling,
        "current_floor": args.current_floor,
        "goodwill_ceiling": args.goodwill_ceiling,
        "neg_ocf_periods": DEFAULT_THRESHOLDS["neg_ocf_periods"],
        "ap_threshold": DEFAULT_THRESHOLDS["ap_threshold"],
    }

    changes = compute_all_changes(data, periods, yoy_map)
    ratios = compute_ratios(data, periods)
    alerts = detect_anomalies(data, periods, changes, ratios, thresholds)

    if args.json or (args.output and args.output.endswith(".json")):
        report = build_json_report(data, periods, changes, ratios, alerts)
        output_text = json.dumps(report, ensure_ascii=False, indent=2) + "\n"
    else:
        output_text = format_text_report(data, periods, changes, ratios, alerts) + "\n"

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(output_text)
        print(f"Report exported to {args.output}", file=sys.stderr)
    else:
        sys.stdout.write(output_text)


if __name__ == "__main__":
    main()
```

## Input Data Format

```json
{
  "company": "Acme Corp",
  "currency": "USD",
  "unit": "thousands",
  "periods": ["2023Q1","2023Q2","2023Q3","2023Q4","2024Q1","2024Q2","2024Q3","2024Q4"],
  "income_statement": {
    "revenue": [5000, 5200, 4800, 6000, 5500, 5800, 5100, 6500],
    "cost_of_revenue": [3000, 3100, 2900, 3500, 3400, 3600, 3200, 4100],
    "operating_income": [800, 850, 750, 1000, 780, 820, 700, 900],
    "net_income": [600, 650, 560, 780, 580, 620, 520, 680]
  },
  "balance_sheet": {
    "accounts_receivable": [2000, 2100, 2200, 2300, 2800, 3200, 3600, 4200],
    "inventory": [1000, 1050, 1100, 1200, 1100, 1150, 1200, 1300],
    "total_current_assets": [5000, 5200, 5400, 5800, 6000, 6500, 7000, 7500],
    "goodwill": [500, 500, 500, 500, 500, 500, 500, 500],
    "total_assets": [15000, 15500, 16000, 16500, 17000, 17500, 18000, 18500],
    "accounts_payable": [1500, 1600, 1550, 1700, 1650, 1750, 1700, 1800],
    "total_current_liabilities": [4000, 4200, 4100, 4500, 4300, 4600, 4500, 4900],
    "total_liabilities": [8000, 8200, 8400, 8600, 8800, 9000, 9200, 9500],
    "total_equity": [7000, 7300, 7600, 7900, 8200, 8500, 8800, 9000]
  },
  "cash_flow": {
    "operating_cash_flow": [700, 750, 620, 850, 300, 280, 250, 200],
    "investing_cash_flow": [-200, -180, -250, -300, -400, -350, -300, -280],
    "financing_cash_flow": [-100, -50, -80, -120, 200, 150, 100, 50],
    "capex": [180, 160, 230, 280, 380, 330, 280, 260]
  }
}
```

`periods` supports quarterly (`2024Q1`) and annual (`2024`) format, auto-detected.
All arrays must match the length of `periods`. Missing fields are skipped
gracefully (no errors thrown). `unit` is a display label only.

## Quick Start

```bash
python3 /workspace/.claw_skills/financial-statement-analyzer/analyze_financials.py data.json
python3 /workspace/.claw_skills/financial-statement-analyzer/analyze_financials.py data.json --json
python3 /workspace/.claw_skills/financial-statement-analyzer/analyze_financials.py data.json -o report.json
python3 /workspace/.claw_skills/financial-statement-analyzer/analyze_financials.py --sample > sample_data.json
```

Key flags: `--json` (JSON output instead of text), `-o/--output <path>` (write
to a file), `--ar-threshold`/`--inv-threshold`/`--ocf-ratio`/`--margin-threshold`/
`--debt-ceiling`/`--current-floor`/`--goodwill-ceiling` (tune anomaly
sensitivity, all default to the values baked into the script).

## Anomaly Detection Rules

| Rule | Trigger | Risk Implication |
|---|---|---|
| AR Surge | AR growth − revenue growth > threshold | Aggressive revenue recognition or collection difficulties |
| Cash Flow Divergence | OCF / Net Income < threshold | Low earnings quality; profits may contain significant accruals |
| Inventory Buildup | Inventory growth − revenue growth > threshold | Product obsolescence or write-down risk |
| Gross Margin Shift | Gross margin change > threshold | Change in pricing power or cost structure |
| Net Margin Shift | Net margin change > threshold | Abnormal expense control or non-recurring items |
| Persistent Negative OCF | OCF < 0 for 2+ consecutive periods | Insufficient organic cash generation |
| Excessive Goodwill | Goodwill / Total Assets > threshold | Impairment risk if acquisitions underperform |
| High Leverage | Liabilities / Assets > threshold | Elevated debt repayment pressure |
| Low Current Ratio | Current Assets / Current Liabilities < threshold | Weak short-term liquidity |
| AP Anomaly | AP growth deviates sharply from COGS growth | Supply chain stress or working-capital strain |

## LLM Interpretation Guide

When a user provides financial data (PDF/image/table/text):
1. **Data Extraction** — convert it into the JSON format above and save it to
   `/workspace/data.json` (or similar).
2. **Run Analysis** — invoke the script above via `exec` for the quantitative
   numbers (don't recompute ratios/YoY/QoQ by hand — trust the script's output).
3. **Comprehensive Interpretation** — combine the script's output with the
   framework below into the final report.

### Cross-Statement Analysis Framework
- **Income Statement → Balance Sheet**: is revenue growth driven by
  receivables? Is net income converting into retained earnings?
- **Income Statement → Cash Flow Statement**: does net income track operating
  cash flow? Are D&A add-backs reasonable?
- **Balance Sheet → Cash Flow Statement**: what's funding asset expansion?
  Is investing activity consistent with capex?

### Contextual Judgment for Anomaly Signals
An anomaly is a signal to investigate, not a verdict — interpret it against
industry/business context: an AR surge can be normal year-end seasonality for
B2B; high leverage is typical for utilities/real estate; negative short-term
cash flow can be reasonable for high-growth SaaS.

### Recommended Output Format
```
## Financial Statement Analysis Report — [Company Name]

### Key Metrics at a Glance
(summary table of key indicators)

### YoY/QoQ Change Highlights
(top 3-5 most significant changes with interpretation)

### Anomaly Signals
(each detected anomaly explained with severity and possible causes)

### Cross-Statement Analysis
(cross-statement logical validation conclusions)

### Summary & Recommendations
(1-2 paragraph overall assessment)
```
"""

_LEGAL_RISK_CONTENT = """\
You are a legal risk assessment assistant for an in-house legal team. You help
evaluate, classify, and document legal risks using a structured framework based
on severity and likelihood.

**Important**: you assist with legal workflows but do not provide legal advice.
Every assessment must be reviewed by qualified legal professionals before it is
relied on. The framework below is a starting point — treat organization-specific
thresholds (financial percentages, escalation names) as illustrative defaults to
confirm with the user, not fixed rules.

## Output Language
Write the final assessment, memo, or risk register entry in Thai (ภาษาไทย) by
default. Only use a different language if the user has explicitly asked for
one (in this request or earlier in the conversation). Keep standard legal/risk
terms of art (e.g. proper names, statute/regulation titles, defined contract
terms) in their original language where a Thai translation would be ambiguous.

## Risk Assessment Framework

**Severity** (impact if the risk materializes), 1-5:
1. Negligible — no material financial/operational/reputational impact.
2. Low — minor exposure (< 1% of relevant value); no public attention.
3. Moderate — material exposure (1-5% of value); noticeable disruption; limited public attention possible.
4. High — substantial exposure (5-25% of value); significant disruption; likely public attention and regulatory scrutiny.
5. Critical — major exposure (> 25% of value); fundamental business disruption; regulatory action likely; possible personal liability for officers/directors.

**Likelihood** (probability it materializes), 1-5:
1. Remote — would require exceptional circumstances, no precedent.
2. Unlikely — possible but not expected, limited precedent.
3. Possible — some precedent, foreseeable triggering events.
4. Likely — clear precedent, triggering events common in similar situations.
5. Almost Certain — strong precedent/pattern, triggering events present or imminent.

**Risk Score = Severity x Likelihood** (1-25):
| Score | Level | Color |
|---|---|---|
| 1-4 | Low Risk | GREEN |
| 5-9 | Medium Risk | YELLOW |
| 10-15 | High Risk | ORANGE |
| 16-25 | Critical Risk | RED |

## Risk Classification Levels

**GREEN — Low (1-4)**: standard business risk. Accept, document in the risk
register, monitor periodically (quarterly/annually), no escalation needed.
E.g. minor deviation from standard vendor terms, routine NDA with a known
counterparty, a routine compliance task with a clear owner/deadline.

**YELLOW — Medium (5-9)**: warrants attention but not immediate action.
Mitigate (negotiate/add controls), monitor at regular intervals (monthly or
on trigger events), document fully in the risk register, assign an owner,
brief stakeholders, and define what would elevate it to ORANGE.
E.g. a below-standard-but-negotiable liability cap, a vendor processing data
in a jurisdiction without a clear adequacy determination, a broader-than-preferred
but market-standard IP clause.

**ORANGE — High (10-15)**: significant issue with meaningful probability.
Escalate to senior counsel/head of legal, brief business leadership, build a
specific mitigation plan with a contingency for if it materializes, consider
outside counsel, review weekly or at milestones, write a full risk memo.
E.g. uncapped indemnification in a material area, threatened litigation from
a significant counterparty, a colorable IP infringement allegation, a
regulatory inquiry/audit request.

**RED — Critical (16-25)**: likely or certain to materialize, could
fundamentally impact the business. Escalate immediately to General
Counsel/C-suite/Board as appropriate, engage outside counsel now, stand up a
response team, consider insurer notification, preserve evidence (litigation
hold) if proceedings are possible, review daily until resolved, make any
required regulatory notifications.
E.g. active litigation with significant exposure, a data breach affecting
regulated personal data, a regulatory enforcement action or government
investigation, a material breach of contract by or against the organization.

## Risk Assessment Memo

Produce this structure for any formal assessment (mark privileged if
applicable):
1. Risk Description
2. Background and Context
3. Risk Analysis — Severity [1-5, label] with rationale; Likelihood [1-5,
   label] with rationale; Risk Score and color
4. Contributing Factors
5. Mitigating Factors
6. Mitigation Options (table: option / effectiveness / cost-effort / recommended?)
7. Recommended Approach
8. Residual Risk (expected level after mitigation)
9. Monitoring Plan (cadence, trigger events for re-assessment)
10. Next Steps (action — owner — deadline)

## Risk Register Entry

For portfolio tracking, capture: Risk ID, Date Identified, Description,
Category (Contract / Regulatory / Litigation / IP / Data Privacy /
Employment / Corporate / Other), Severity, Likelihood, Risk Score, Risk
Level, Owner, Mitigations, Status (Open / Mitigated / Accepted / Closed),
Review Date, Notes.

## When to Escalate to Outside Counsel

**Mandatory**: active litigation, a government/regulator investigation,
potential criminal exposure, anything touching securities disclosures/filings,
any matter requiring board notification or approval.

**Strongly recommended**: novel/first-impression legal issues, unfamiliar or
conflicting jurisdictions, exposure beyond the organization's risk tolerance,
specialized expertise not available in-house (antitrust, FCPA, patent
prosecution, etc.), new regulation requiring a compliance program, M&A
due diligence/deal structuring/regulatory approvals.

**Consider**: significant contract interpretation disputes, employment
claims (discrimination, harassment, wrongful termination, whistleblower),
a potential data breach with notification obligations, IP disputes over
material products, insurance coverage disputes.

When recommending engagement, note the factors the user should weigh in
selecting counsel: subject-matter expertise, jurisdiction experience,
industry familiarity, conflict clearance, fee arrangement, and any existing
panel-firm relationship.
"""

_BUILTIN_SKILLS: tuple[BuiltinSkill, ...] = (
    BuiltinSkill(
        name="skill-creator",
        description=(
            "How to author a new reusable skill correctly — use whenever the user asks you to "
            "create, build, or save a skill."
        ),
        content=_SKILL_CREATOR_CONTENT,
    ),
    BuiltinSkill(
        name="pptx",
        description=(
            "Create and edit PowerPoint presentations — HTML-previewed layouts, template reuse, "
            "charts/tables, and slide analysis. Use whenever the user asks for a .pptx/slide deck."
        ),
        content=_PPTX_CONTENT,
        capabilities=(
            (
                "Presentation Creation",
                "HTML-to-PPTX workflow with custom color palettes, charts, tables, and visual design",
            ),
            (
                "Template Reuse",
                "automatically matching layouts, duplicating slides, and replacing content",
            ),
            (
                "Slide Editing",
                "OOXML manipulation, modifying text, layouts, and design elements",
            ),
            (
                "Thumbnails & Analysis",
                "slide thumbnail grids, extract text content, and analyze layout structure",
            ),
        ),
        summary=(
            "Build PowerPoint presentations for supplier reviews, data summaries, and strategy "
            "proposals with structured layouts and chart integration."
        ),
    ),
    BuiltinSkill(
        name="xlsx",
        description=(
            "Create, analyze, and visualize Excel spreadsheets — themed tables/charts, pandas "
            "analysis, and financial models. Use whenever the user asks for a .xlsx/spreadsheet."
        ),
        content=_XLSX_CONTENT,
        capabilities=(
            (
                "Spreadsheet Creation",
                "build_table/build_chart with multiple theme styles and automatic formula calculation",
            ),
            (
                "Data Analysis",
                "pandas processing, statistical calculations, and structured output",
            ),
            (
                "Chart Visualization",
                "bar, line, area, and combo charts with themed colors, axis configuration, and legend layouts",
            ),
            (
                "Financial Models",
                "DCF, LBO, three-statement models with formula linking, source annotations, and scenario analysis",
            ),
        ),
        summary=(
            "Create, edit, analyze, and visualize Excel spreadsheets with formulas, formatting, "
            "charts, and data cleaning — from supplier comparison tables to financial models."
        ),
    ),
    BuiltinSkill(
        name="pdf",
        description=(
            "Extract, create, merge/split, and fill PDFs — layout-preserving extraction, branded "
            "document generation, and AcroForm filling. Use whenever the user asks for a .pdf."
        ),
        content=_PDF_CONTENT,
        capabilities=(
            (
                "Text & Table Extraction",
                "layout preservation and structured output",
            ),
            (
                "PDF Creation",
                "using reportlab or PyMuPDF, supporting multi-page, charts, and custom styling",
            ),
            (
                "Merge & Split",
                "supporting page rotation and reordering",
            ),
            (
                "Form Filling",
                "field detection, auto-fill, and validation",
            ),
        ),
        summary=(
            "Generate branded PDF documents, extract key clauses from contracts, and merge or "
            "split PDF files for procurement and business use cases."
        ),
    ),
    BuiltinSkill(
        name="docx",
        description=(
            "Create and edit Word documents — professional typesetting, tracked changes/review "
            "comments, and content extraction. Use whenever the user asks for a .docx/Word document."
        ),
        content=_DOCX_CONTENT,
        capabilities=(
            (
                "Document Creation",
                "using python-docx, supporting paragraphs, tables, headers/footers, and professional typesetting",
            ),
            (
                "Document Editing",
                "OOXML manipulation, preserving original formatting, styles, and formulas",
            ),
            (
                "Tracked Changes",
                "apply tracked changes and review comments, suitable for legal, academic, and business document review",
            ),
            (
                "Content Extraction",
                "document text, metadata, embedded images, and comments, supporting multiple output formats",
            ),
        ),
        summary=(
            "Create and edit Word documents including RFQ letters, supplier evaluation reports, "
            "SOPs, and contract templates with professional formatting."
        ),
    ),
    BuiltinSkill(
        name="legal-risk-assessment",
        description=(
            "Assess and classify legal risks using a severity x likelihood framework with "
            "escalation criteria. Use when evaluating contract risk, assessing deal exposure, "
            "classifying issues by severity, or deciding whether a matter needs senior counsel "
            "or outside legal review."
        ),
        content=_LEGAL_RISK_CONTENT,
        capabilities=(
            (
                "Risk Scoring",
                "severity x likelihood matrix (1-5 each) producing a 1-25 score mapped to GREEN/YELLOW/ORANGE/RED risk levels",
            ),
            (
                "Escalation Criteria",
                "clear thresholds for when a matter needs senior counsel, outside counsel, or board/executive attention",
            ),
            (
                "Risk Assessment Memo",
                "structured 10-section memo — description, severity/likelihood rationale, mitigation options, residual risk, monitoring plan",
            ),
            (
                "Risk Register Tracking",
                "standardized fields (category, owner, status, review date) for ongoing risk portfolio tracking",
            ),
        ),
        summary=(
            "Classify legal risks by severity and likelihood, decide whether to escalate to "
            "senior or outside counsel, and produce a structured risk assessment memo or risk "
            "register entry. Not a substitute for review by qualified legal professionals."
        ),
    ),
    BuiltinSkill(
        name="financial-statement-analyzer",
        description=(
            "Analyzes income statement, balance sheet, and cash flow statement data to generate "
            "YoY/QoQ trend analysis and flag anomalies like AR surges or cash flow divergence. "
            "Use when asked to analyze financials, compare YoY/QoQ, detect red flags, or assess "
            "earnings quality."
        ),
        content=_FINANCIAL_STATEMENT_CONTENT,
        capabilities=(
            (
                "YoY Analysis",
                "compare same-period data (e.g. 2024Q1 vs 2023Q1) to identify trend changes",
            ),
            (
                "QoQ Analysis",
                "compare consecutive periods (e.g. 2024Q2 vs 2024Q1) to capture short-term fluctuations",
            ),
            (
                "Financial Ratios",
                "gross margin, net margin, debt-to-asset ratio, current ratio, DSO, and more",
            ),
            (
                "Anomaly Detection",
                "10 built-in rules with automatic scanning, risk severity levels, and explanations",
            ),
        ),
        summary=(
            "Compute YoY/QoQ trends and financial ratios from income statement, balance sheet, "
            "and cash flow data, then run 10 anomaly detection rules (AR surge, cash flow "
            "divergence, inventory buildup, and more) to flag earnings-quality red flags."
        ),
    ),
)


def builtin_skills() -> list[BuiltinSkill]:
    """All built-in skills (always enabled)."""
    return list(_BUILTIN_SKILLS)


def get_builtin_skill(name: str) -> BuiltinSkill | None:
    for skill in _BUILTIN_SKILLS:
        if skill.name == name:
            return skill
    return None
